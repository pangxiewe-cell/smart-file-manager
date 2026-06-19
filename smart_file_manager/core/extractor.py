"""
文件内容提取器
根据文件类型选用不同库提取文本内容，不依赖 Java/tika。
"""
from pathlib import Path
import config


MAX_CHARS = config.LLM_MAX_TEXT_CHARS


def extract_text(path: str) -> str:
    """根据文件类型提取文本，失败时返回空字符串"""
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext in {".txt", ".md", ".py", ".js", ".ts", ".json",
                   ".yaml", ".yml", ".csv", ".html", ".css"}:
            return _read_plain(p)
        elif ext == ".pdf":
            return _read_pdf(p)
        elif ext == ".docx":
            return _read_docx(p)
        elif ext == ".xlsx":
            return _read_xlsx(p)
        elif ext == ".pptx":
            return _read_pptx(p)
        elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}:
            return _read_image_ocr(p)
    except Exception as e:
        print(f"[Extractor] 提取失败 {p.name}：{e}")
    return ""


def _read_plain(p: Path) -> str:
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            return p.read_text(encoding=enc, errors="replace")[:MAX_CHARS]
        except Exception:
            continue
    return ""


def _read_pdf(p: Path) -> str:
    try:
        from pdfminer.high_level import extract_text as pdf_extract
        return (pdf_extract(str(p)) or "")[:MAX_CHARS]
    except Exception:
        # pdfminer 失败，降级用 pypdfium2
        try:
            import pypdfium2 as pdfium
            doc = pdfium.PdfDocument(str(p))
            parts = []
            for i in range(min(len(doc), 20)):
                page = doc[i]
                parts.append(page.get_textpage().get_text_range())
            return "\n".join(parts)[:MAX_CHARS]
        except Exception:
            return ""


def _read_docx(p: Path) -> str:
    from docx import Document
    doc = Document(str(p))
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    return text[:MAX_CHARS]


def _read_xlsx(p: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(p), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets[:3]:
        for row in ws.iter_rows(max_row=50, values_only=True):
            cell_text = " ".join(str(c) for c in row if c is not None)
            if cell_text.strip():
                lines.append(cell_text)
    return "\n".join(lines)[:MAX_CHARS]


def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    texts = []
    for slide in prs.slides[:10]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    texts.append(txt)
    return "\n".join(texts)[:MAX_CHARS]


def _read_image_ocr(p: Path) -> str:
    """OCR 提取图片文字（仅图片小于 400 万像素时执行）"""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(str(p))
        if img.width * img.height > 4_000_000:
            return ""
        return pytesseract.image_to_string(img, lang="chi_sim+eng")[:MAX_CHARS]
    except Exception:
        return ""


def compute_content_hash(path: str, extracted_text: str) -> str:
    """
    计算内容哈希：
    - 文本类文件（有足够文本）→ SimHash（可检测近重复）
    - 其他 → xxhash（速度快，碰撞率低）
    """
    ext = Path(path).suffix.lower()
    text_exts = {".txt", ".md", ".py", ".js", ".ts", ".docx", ".pdf", ".csv"}
    if ext in text_exts and len(extracted_text) > 100:
        try:
            from simhash import Simhash
            return f"sh:{Simhash(extracted_text).value}"
        except Exception:
            pass
    # 二进制类：只读前 1MB 计算快速哈希
    try:
        import xxhash
        with open(path, "rb") as f:
            data = f.read(1 * 1024 * 1024)
        return f"xx:{xxhash.xxh64(data).hexdigest()}"
    except Exception:
        return ""
